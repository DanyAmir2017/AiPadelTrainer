from pptx import Presentation
from pptx.util import Pt
from pathlib import Path

out_path = Path("DOCTOR_7_SLIDES.pptx")

slides_data = [
    {
        "title": "Padel Ball Tracking Thesis: Problem, Objective, Scope",
        "bullets": [
            "Challenge: tiny fast ball, motion blur, perspective distortion, occlusion",
            "Objective: detect/track ball, evaluate shot zones, export analytics",
            "Scope delivered: robust CV pipeline plus edge deployment pathway",
            "Detection chain: YOLO -> Optical Flow -> Kalman -> outputs",
        ],
        "notes": [
            "Start with why this is difficult in real padel footage.",
            "Explain this is both an AI and systems-engineering thesis.",
            "Stress practical usability for training feedback.",
            "Transition to architecture that enabled iteration.",
        ],
    },
    {
        "title": "System Architecture and Technical Stack",
        "bullets": [
            "Modules: detection, tracking, filtering, evaluation, visualization/output",
            "Stack: Python, Ultralytics YOLOv8, OpenCV, FilterPy, ONNX Runtime",
            "Models: best_ball, best_players, best_court with ONNX exports",
            "Fallback strategy improved robustness when primary detector missed",
        ],
        "notes": [
            "Show modularity enabled controlled experiments and targeted fixes.",
            "Mention ONNX as bridge toward deployment portability.",
            "Clarify fallback chain avoided brittle single-model behavior.",
            "Prepare audience for baseline bottlenecks next.",
        ],
    },
    {
        "title": "Baseline Results and Bottlenecks",
        "bullets": [
            "640px baseline rates: V1 40.6%, V2 74.0%, V3 15.1%, V5 39.6%, V6 24.4%",
            "Typical source split: YOLO dominant, optical flow critical fallback",
            "Video 1 showed fallback strength (Optical Flow 53.1%)",
            "Main bottlenecks: blur, tiny object scale, spatial bias (bottom-left weakness)",
        ],
        "notes": [
            "Frame this as diagnostic baseline, not final performance.",
            "Point out high variance across videos and contexts.",
            "Explain why these issues required multi-pronged improvements.",
            "Transition into experiment log and decisions.",
        ],
    },
    {
        "title": "Experiment Log: What Worked vs What Failed",
        "bullets": [
            "Fine-tuning attempt: no meaningful improvement -> not deployed",
            "CLAHE on Video 5: 237->227 detections and 10.2->7.9 FPS -> rejected",
            "Region-adaptive confidence: no bottom-left recovery -> rejected",
            "1280 inference: better coverage (+37 bottom-left on V6) but major speed trade-off",
        ],
        "notes": [
            "Highlight evidence-based decision making.",
            "Negative results reduced wasted complexity in final system.",
            "Keep emphasis on reproducible measurements.",
            "Lead into edge deployment outcomes.",
        ],
    },
    {
        "title": "Edge Deployment Progress (Raspberry Pi 5 + Hailo)",
        "bullets": [
            "Edge-focused refactor completed for deployment constraints",
            "Pi ONNX run validated (~5.3 FPS documented first successful profile)",
            "Critical finding: ONNX->HEF compile requires x86_64, unavailable on ARM64 Pi",
            "Hailo runtime verified with precompiled model at ~154.96 FPS",
        ],
        "notes": [
            "Differentiate runtime success from compiler/toolchain limitation.",
            "Show this strengthens deployment realism in thesis.",
            "Explain practical workaround paths (x86_64 compile route).",
            "Transition to latest contact-validation milestone.",
        ],
    },
    {
        "title": "Contact Validation Milestone and Recent Gains",
        "bullets": [
            "Added clean trajectory + rule-based contact candidate extraction",
            "Built labeling GUI and added contact types (racket/ground/glass/out_of_frame)",
            "Video 5: v1 72.0% (18/25) -> v2 90.0% (18/20)",
            "Video 6: v2 56.604% (30/53) -> v3 73.684% (28/38)",
            "Latest Pi rerun (v5_after_tune_v3): 5-second snapshots verified + frame 598 remainder",
        ],
        "notes": [
            "This is the strongest recent quality-improvement evidence.",
            "Accuracy increased while reducing false positives.",
            "Call out unspecified-contact ambiguity as the next quality target.",
            "Then summarize total contribution and completion plan.",
        ],
    },
    {
        "title": "Current Status, Contributions, and Next Steps",
        "bullets": [
            "Current status: end-to-end system + experiments + deployment validation complete",
            "Contributions: robust fallback design, quantified trade-offs, deployment-aware engineering",
            "Validated contact-event workflow with measurable accuracy gains",
            "Next: finalize benchmark chapter, optional x86_64 Hailo compile, thesis/demo finalization",
            "Supervisor ask: feedback on final evaluation criteria and demo scope",
        ],
        "notes": [
            "Conclude with confidence and evidence-backed progress.",
            "Reinforce that both research and implementation objectives were achieved.",
            "Present clear, finite remaining milestones.",
            "Invite supervisor feedback on final evaluation framing.",
        ],
    },
]

prs = Presentation()
for idx, item in enumerate(slides_data):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = item["title"]

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    for i, bullet in enumerate(item["bullets"]):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22 if idx == 0 else 20)

    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = "\n".join(item["notes"])

prs.save(out_path)
print(f"Created: {out_path.resolve()}")
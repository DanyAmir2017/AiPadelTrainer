from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


OUT_PATH = Path("padel_trainer_results_tables_v6.pptx")


RAW_DETECTION_ROWS = [
    ["Padel_video_7", "1030", "1028", "688", "99.81%", "12.70", "0.00"],
    ["Padel_video_8", "1431", "1429", "789", "99.86%", "17.15", "0.00"],
]

CONTACT_TUNING_ROWS = [
    ["Video 5 v1", "25", "18", "7", "72.0%"],
    ["Video 5 v2", "20", "18", "2", "90.0%"],
    ["Video 5 improvement", "-5 candidates", "+0 correct", "-5 wrong", "+18.0 pts"],
    ["Video 6 V2", "53 labeled", "30 correct", "23 wrong", "56.604%"],
    ["Video 6 V3", "38 labeled", "28 correct", "10 wrong", "73.684%"],
]

VIDEO6_SCORER_ROWS = [
    ["TP", "21", "30", "+9"],
    ["FP", "17", "9", "-8"],
    ["FN", "9", "0", "-9"],
    ["Precision", "55.263%", "76.923%", "+21.660 pts"],
    ["Recall", "70.000%", "100.000%", "+30.000 pts"],
    ["F1", "61.765%", "86.957%", "+25.192 pts"],
    ["Accuracy", "44.681%", "76.923%", "+32.242 pts"],
]

PI_LAPTOP_ROWS = [
    ["Video 7 - Trajectory rows", "1030", "1030", "0"],
    ["Video 7 - Clean trajectory rows", "761", "688", "+73"],
    ["Video 7 - Hit candidates", "42", "34", "+8"],
    ["Video 7 - Accept / Review / Reject", "36 / 4 / 2", "23 / 9 / 2", "+13 / -5 / 0"],
    ["Video 8 - Trajectory rows", "1431", "1431", "0"],
    ["Video 8 - Clean trajectory rows", "855", "789", "+66"],
    ["Video 8 - Hit candidates", "39", "26", "+13"],
    ["Video 8 - Accept / Review / Reject", "32 / 4 / 3", "23 / 3 / 0", "+9 / +1 / +3"],
]

MANUAL_VS_MODEL_ROWS = [
    ["Video 6", "46", "38", "-8"],
    ["Video 7", "36", "42", "+6"],
    ["Video 8", "44", "39", "-5"],
]

ADDITIONAL_METRICS_ROWS_1 = [
    ["Video 3 fine-tune", "108 frames", "11 epochs", "mAP50 = 0.479", "Not deployed"],
]

ADDITIONAL_METRICS_ROWS_2 = [
    ["Training samples", "156"],
    ["Accept threshold", "0.5711"],
    ["Review threshold", "0.4011"],
    ["Metrics @ accept", "acc=0.936, prec=0.945, rec=0.963, f1=0.954"],
    ["Output model", "windowed_contact_scorer_v2_v5_v8.json"],
]


def set_bg(slide, color=(255, 255, 255)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.20), Inches(12.45), Inches(0.55))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(17, 24, 39)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.47), Inches(0.72), Inches(12.2), Inches(0.35))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.size = Pt(10.5)
        srun.font.color.rgb = RGBColor(75, 85, 99)


def style_cell(cell, font_size=11, bold=False, align=PP_ALIGN.CENTER, fill=None, text_color=(17, 24, 39)):
    tf = cell.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.bold = bold
            r.font.color.rgb = RGBColor(*text_color)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*fill)


def set_table_text(cell, text, font_size=11, bold=False, align=PP_ALIGN.CENTER):
    cell.text = str(text)
    style_cell(cell, font_size=font_size, bold=bold, align=align)


def add_table(slide, left, top, width, height, headers, rows, col_widths=None, font_size=11):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)

    header_fill = (30, 58, 138)
    body_fill = (248, 250, 252)
    alt_fill = (241, 245, 249)
    grid_color = RGBColor(203, 213, 225)

    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        set_table_text(cell, head, font_size=font_size, bold=True, align=PP_ALIGN.CENTER)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*header_fill)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(255, 255, 255)

    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            set_table_text(cell, value, font_size=font_size, bold=False, align=align)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*(body_fill if i % 2 else alt_fill))

    for row in table.rows:
        for cell in row.cells:
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.text_frame.word_wrap = True
            for p in cell.text_frame.paragraphs:
                p.space_after = Pt(0)
            # borders
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
                ln = tcPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}lnL") if edge == "a:lnL" else None
            # Borders are controlled by the theme in PowerPoint; keep the clean fill and header contrast.

    return table_shape


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.00), Inches(12.3), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.color.rgb = RGBColor(107, 114, 128)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(0.9))
    tf = title.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Padel Trainer Thesis Results"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(17, 24, 39)
    sub = slide.shapes.add_textbox(Inches(0.72), Inches(2.15), Inches(11.8), Inches(0.5))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = "Detection coverage, contact validation, and deployment comparison tables"
    sr.font.size = Pt(14)
    sr.font.color.rgb = RGBColor(75, 85, 99)
    foot = slide.shapes.add_textbox(Inches(0.72), Inches(6.5), Inches(11.8), Inches(0.4))
    f = foot.text_frame.paragraphs[0].add_run()
    f.text = "Generated from padel_trainer outputs and documentation"
    f.font.size = Pt(10.5)
    f.font.color.rgb = RGBColor(107, 114, 128)

    # Slide 2: raw detection results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Raw Ball Detection Coverage Across Videos 7–8", "Detection rate = total ball detections / total frames processed")
    headers = ["Video", "Frames", "Raw", "Clean", "Detection rate", "FPS", "Consistency"]
    add_table(
        slide,
        Inches(0.45), Inches(1.15), Inches(12.4), Inches(5.7),
        headers,
        RAW_DETECTION_ROWS,
        col_widths=[2.2, 1.0, 1.0, 1.0, 1.4, 1.0, 1.2],
        font_size=11,
    )
    add_footer(slide, "Source: outputs/metrics/Padel_video_*_metrics.txt")

    # Slide 3: contact tuning
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Contact Validation Tuning Results", "Evidence-based refinement of rule-based contact candidates")
    headers = ["Run", "Candidates", "Correct", "Wrong", "Accuracy"]
    add_table(
        slide,
        Inches(0.45), Inches(1.15), Inches(12.4), Inches(5.7),
        headers,
        CONTACT_TUNING_ROWS,
        col_widths=[3.2, 1.6, 1.6, 1.4, 1.5],
        font_size=11,
    )
    add_footer(slide, "Source: DOCUMENTATION.md Section 13.11.4–13.11.5 and labeled candidate CSVs")

    # Slide 4: scorer metrics
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Video 6 Contact Scorer Performance", "Learned scorer vs rule-only baseline")
    headers = ["Metric", "V3 rule-only", "Scored-V4", "Delta"]
    add_table(
        slide,
        Inches(0.45), Inches(1.15), Inches(12.4), Inches(5.7),
        headers,
        VIDEO6_SCORER_ROWS,
        col_widths=[2.6, 2.0, 2.0, 2.4],
        font_size=11,
    )
    add_footer(slide, "Source: DOCUMENTATION.md Section 13.11.8 and 13.11.9")

    # Slide 5: Pi vs Laptop comparison
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Video 7/8 Pi vs Laptop Comparison", "Structural parity with expected runtime variability")
    headers = ["Metric", "Laptop", "Pi", "Delta"]
    add_table(
        slide,
        Inches(0.35), Inches(1.08), Inches(12.65), Inches(5.85),
        headers,
        PI_LAPTOP_ROWS,
        col_widths=[4.1, 1.7, 1.7, 2.1],
        font_size=10.5,
    )
    add_footer(slide, "Source: DOCUMENTATION.md Section 13.11.15–13.11.16")

    # Slide 6: manual vs model collision counts
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Manual Contact Counts vs Model Collision Detections", "Model values are the final collision-candidate totals across ground, glass, and racket contacts")
    headers = ["Video", "Manual count", "Model collisions", "Delta"]
    add_table(
        slide,
        Inches(1.1), Inches(1.35), Inches(11.1), Inches(3.4),
        headers,
        MANUAL_VS_MODEL_ROWS,
        col_widths=[2.2, 2.2, 2.3, 2.2],
        font_size=12,
    )
    add_footer(slide, "Source: manual counts provided by user; model counts from documentation and hit-candidate outputs")

    # Slide 7: additional metrics
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Additional Training and Benchmark Metrics", "Useful supporting results for the thesis evaluation chapter")

    add_table(
        slide,
        Inches(0.45), Inches(1.15), Inches(12.4), Inches(1.35),
        ["Fine-tuning run", "Frames", "Epochs", "Metric", "Decision"],
        ADDITIONAL_METRICS_ROWS_1,
        col_widths=[3.4, 1.4, 1.1, 2.0, 2.3],
        font_size=11,
    )

    add_table(
        slide,
        Inches(0.45), Inches(2.95), Inches(12.4), Inches(3.25),
        ["Windowed scorer summary", "Value"],
        ADDITIONAL_METRICS_ROWS_2,
        col_widths=[3.3, 9.0],
        font_size=11,
    )
    add_footer(slide, "Sources: DOCUMENTATION.md Section 13.11.1, 13.11.12, and 13.11.13")

    return prs


if __name__ == "__main__":
    out_path = Path(__file__).resolve().parents[1] / OUT_PATH
    deck = build_deck()
    deck.save(out_path)
    print(f"Created: {out_path}")
